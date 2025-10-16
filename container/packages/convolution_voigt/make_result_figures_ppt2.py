# ------------------------------------------------------------
# Copyright (c) 2024, National Institute for Materials Science
#
# This software is released under the MIT License.
#
# Contributor:
#     Hiroshi Shinotsuka
# ------------------------------------------------------------
# coding: utf-8
import numpy as np
import pandas as pd
from glob import glob
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import os
import sys
import re

# import os, sys; sys.path.append(os.path.expanduser("~")); from mytemplate import init_figure, init_pptx, pptx_layout_figures


# ------------------------------------------------------------------
def init_pptx():
    """
    Create a new Presentation() instance.
    The slide height is adjusted to make a layout 16:9.
    """
    # Presentationインスタンスの作成。16:9サイズにするためにslide_height = 5143500とする。標準は6858000。
    ppt = Presentation()
    ppt.slide_height = int(ppt.slide_width / 16 * 9)
    return ppt


# ------------------------------------------------------------------
def pptx_layout_figures(ppt, figPathList, nx=None, title=None, spacer=0):
    """
    Add a new slide to ppt, and put figures from figPathList in a grid.
    nx is the number of figures in a row.
    ny, the number of figures in a column is automatically determined (can not be specified).
    When nx = None, the suitable nx and ny will be automatically determined.
    spacer is a space distance between figures.
    title is displayed at the top of the slide.
    """
    # 幅
    width = ppt.slide_width
    # 高さ
    height = ppt.slide_height

    # 画像ファイルの幅と高さを取得。すべての画像の大きさが同じと仮定する。
    # im = Image.open(figPathList[0])

    # print(f"figPathList = ", figPathList)
    for figPath in figPathList:
        if os.path.exists(figPath):
            im = Image.open(figPath)

    imWidth, imHeight = im.size

    # 1行あたりnx個の図を並べるようにレイアウトを設定する
    tH = int(height * 0.08)  # text height
    pT = int(tH * 1.2)  # picture Top position
    pL = int(width * 0.00)  # picture Left position

    numFigs = len(figPathList)
    if nx is None:
        # nxを自動的に決める。width * (height-pT) の枠に全ての画像がはまるように。
        print("nx is automatically determined.")
        nx = 0
        while nx < numFigs:
            nx += 1
            # pW = int(width / nx)   # picture width
            pW = int((width - (nx - 1) * spacer) / nx)  # picture width
            pH = int(pW / imWidth * imHeight)  # picture Height
            ny = int(np.ceil(numFigs / nx))
            # print(nx, ny, pW, pH, (height-pT), ny * pH)
            # if (height-pT) >= ny * pH:
            #    break
            if (height - pT) >= ny * pH + (ny - 1) * spacer:
                break
    else:
        # pW = int(width / nx)   # picture width
        pW = int((width - (nx - 1) * spacer) / nx)  # picture width
        pH = int(pW / imWidth * imHeight)  # picture Height
        ny = int(np.ceil(numFigs / nx))

    if (height - pT) < ny * pH + (ny - 1) * spacer:
        # y方向にはみ出る場合，画像のサイズを小さくする
        # pH = int(np.ceil((height-pT) / ny))
        pH = int(np.ceil((height - pT - (ny - 1) * spacer) / ny))
        pW = int(pH / imHeight * imWidth)
        pL = int((width - nx * pW - (nx - 1) * spacer) / 2)

    ny = int(np.ceil(numFigs / nx))
    print(f"nx={nx}\nny={ny}")

    # pW = int(width / nx)   # picture width
    # pH = int(pW / imWidth * imHeight)  # picture Height

    # 新規スライド。レイアウト6番＝白紙。
    # slide = ppt.slides.add_slide(blank_slide_layout)
    slide = ppt.slides.add_slide(ppt.slide_layouts[6])

    # テキスト
    if title is not None:
        txBox = slide.shapes.add_textbox(int(width * 0), 0.0, int(width), int(height * 0.08))
        tf = txBox.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(20)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ファイルごとにループ
    numFigs = len(figPathList)
    cnt = -1
    for figPath in figPathList:
        cnt += 1
        print("%d/%d, Directory %s" % (cnt + 1, numFigs, figPath))

        # 横10個，縦6個になるように位置を決める
        # pL1 = pL + (cnt % nx) * pW
        # pT1 = pT + int(cnt / nx) * pH
        pL1 = pL + (cnt % nx) * (pW + spacer)
        pT1 = pT + int(cnt / nx) * (pH + spacer)
        # print("%5.1f, %5.1f" % (pL1/width*100, pT1/height*100))

        # pic = slide.shapes.add_picture(figPath, pL1, pT1, pW, pH)

        if os.path.exists(figPath):
            pic = slide.shapes.add_picture(figPath, pL1, pT1, pW, pH)

        # pic.line.color.rgb = RGBColor(0xFF, 0x00, 0xFF)
        # pic.line.width = Inches(0.01)

        continue

    return


def get_fig_vs_numpeak():
    fig_list = glob("*_vs_NumPeak.png")
    if len(fig_list) == 0:
        fig_vs_numpeak = None
    else:
        fig_vs_numpeak = fig_list[0]
    return fig_vs_numpeak


# ------------------------------------------------------------------
def getFigList2():
    figList = []

    # 1st figure is BIC vs Number of peaks
    # figList.append("BIC_vs_NumPeak.png")

    # figListTemp = glob("gbp_rank*_numPeak*_*_mSG*_result.png")
    figListTemp = glob("gbp_rank*_numPeak*_*_result.png")

    if len(figListTemp) == 0:
        print("Error. No result found")
        sys.exit(1)

    figListData = []
    for f in figListTemp:
        # print(f)
        patternInfo = r"gbp_rank(\d+)_numPeak(\d+)_(.*)_result.png"
        match = re.search(patternInfo, f)
        if match is None:
            print("Not match")
            sys.exit(1)
        else:
            rank = match.groups()[0]
            numPeak = match.groups()[1]
            # method = match.groups()[2]
            if "shirley" in match.groups()[2]:
                method = "shirley"
            elif "linear" in match.groups()[2]:
                method = "linear"
            else:
                method = ""
        figListData.append([rank, numPeak, method, f])

    # print(figListData)
    # print(np.array(figListData))
    # print(pd.DataFrame(figListData))
    figListDf = pd.DataFrame(np.array(figListData), columns=["rank", "numPeak", "method", "filename"]).astype({"rank": int, "numPeak": int})
    # figListDf.sort_values('rank', inplace=True)
    # figListDf.sort_values('numPeak', inplace=True)
    figListDf.sort_values(by=["numPeak", "rank"], inplace=True)
    # print(figListDf)

    # print(figList)
    return figListDf


# ------------------------------------------------------------------
def set_text(slide, msg, text_left, text_top, text_width, text_height, fontsize=16, alignment=PP_ALIGN.CENTER):
    txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
    tf = txBox.text_frame
    tf.text = msg
    tf.paragraphs[0].font.size = Pt(fontsize)
    tf.paragraphs[0].alignment = alignment


# ------------------------------------------------------------------

if __name__ == "__main__":

    fig_vs_numpeak = get_fig_vs_numpeak()

    figListDf = getFigList2()
    # print(figListDf)

    # Presentationインスタンスの作成。16:9サイズにするためにslide_height = 5143500とする。標準は6858000。
    ppt = Presentation()
    ppt.slide_height = int(ppt.slide_width / 16 * 9)

    slide = ppt.slides.add_slide(ppt.slide_layouts[6])

    # 幅
    width = ppt.slide_width
    # 高さ
    height = ppt.slide_height

    text2 = "Fitting models selected by BIC, grouped by number of peaks"

    text_left = int(width * 0 / 6.0)
    text_top = int(height * 0.0)
    text_width = int(width * 5.0 / 6.0)
    text_height = int(height * 0.06)
    set_text(slide, text2, text_left, text_top, text_width, text_height, fontsize=16, alignment=PP_ALIGN.LEFT)

    text_left = int(width * 0.0)
    text_top = int(height * 0.06)
    text_width = int(width / 3.0)
    text_height = int(height * 0.06)

    pic_left = int(width * 0.0)
    # pic_top   = int(height * 0.12)
    pic_top = text_top + int(height * 0.02)
    pic_width = int(width / 3.0)
    pic_height = int(pic_width * 3 / 4)

    # textList = ["BIC overview"]
    # textList += [f"Rank {rank}" for rank in figListDf['rank']]

    # figList = ["BIC_vs_NumPeak.png"] + list(figListDf["filename"])
    figList = [fig_vs_numpeak] + list(figListDf["filename"])
    textList = ["BIC overview"]
    for index, row in figListDf.iterrows():
        # textList += [f"Rank {row['rank']}, K={row['numPeak']}"]
        textList += [f"  Rank {row['rank']}, K={row['numPeak']}, {row['method']}"]

    cnt = 0
    for fig, text in zip(figList, textList):
        if cnt == 3:
            pic_left = 0
            text_left = 0
            text_top = pic_top + pic_height
            pic_top = text_top + int(height * 0.02)

        if cnt == 6:
            print("New slide")
            text_left = 0
            text_top = int(height * 0.06)
            pic_left = 0
            pic_top = text_top + int(height * 0.02)
            slide = ppt.slides.add_slide(ppt.slide_layouts[6])
            cnt = 0

        if os.path.exists(fig):
            print(fig)
            slide.shapes.add_picture(fig, pic_left, pic_top, pic_width, pic_height)
            set_text(slide, text, text_left, text_top, text_width, text_height, fontsize=16)

        pic_left += pic_width
        text_left += pic_width

        cnt += 1

    # # 枠
    # shp = slide.shapes.add_textbox(0, int(height * 0.18) + pic_height, width, (text_height + pic_height) * 2)
    # # shp.text = 'sample'
    # shp.line.color.rgb = RGBColor(0, 0, 0)

    pptfile = "result_figures.pptx"
    print("Output to", pptfile)
    ppt.save(pptfile)
